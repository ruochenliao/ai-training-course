"""
基于Marker框架的文档解析服务 - 企业级RAG系统
严格按照技术栈要求：VikParuchuri/marker 文档解析引擎
支持格式：PDF、DOCX、PPTX、XLSX、MD、TXT (最大100MB)
"""
import asyncio
import hashlib
from datetime import datetime
from pathlib import Path
from typing import Dict, Tuple, Any

import aiofiles
from app.core.config import settings
from app.core.database_new import get_db_session
from app.models.sqlalchemy_models import Document, ParseStatus
from loguru import logger
from marker import convert_single_pdf
from marker.models import load_all_models
from marker.settings import settings as marker_settings


class MarkerDocumentService:
    """基于Marker框架的文档解析服务"""
    
    def __init__(self):
        self.models = None
        self.supported_formats = {
            'pdf': 'application/pdf',
            'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            'md': 'text/markdown',
            'txt': 'text/plain',
        }
        self.max_file_size = 100 * 1024 * 1024  # 100MB
        self.upload_dir = Path(settings.UPLOAD_DIR)
        self.upload_dir.mkdir(exist_ok=True)
    
    async def initialize_models(self):
        """初始化Marker模型 - 异步加载"""
        if self.models is None:
            try:
                logger.info("正在加载Marker模型...")
                # 在线程池中加载模型，避免阻塞
                loop = asyncio.get_event_loop()
                self.models = await loop.run_in_executor(None, load_all_models)
                logger.info("Marker模型加载完成")
            except Exception as e:
                logger.error(f"Marker模型加载失败: {e}")
                raise
    
    def validate_file(self, filename: str, file_size: int, content_type: str) -> Tuple[bool, str]:
        """验证文件格式和大小"""
        # 检查文件大小
        if file_size > self.max_file_size:
            return False, f"文件大小超过限制 ({self.max_file_size / 1024 / 1024:.1f}MB)"
        
        # 检查文件扩展名
        file_ext = Path(filename).suffix.lower().lstrip('.')
        if file_ext not in self.supported_formats:
            return False, f"不支持的文件格式: {file_ext}"
        
        # 检查MIME类型
        expected_mime = self.supported_formats[file_ext]
        if content_type != expected_mime:
            logger.warning(f"MIME类型不匹配: 期望 {expected_mime}, 实际 {content_type}")
        
        return True, "文件验证通过"
    
    async def calculate_file_hash(self, file_path: Path) -> str:
        """计算文件MD5哈希值"""
        hash_md5 = hashlib.md5()
        async with aiofiles.open(file_path, 'rb') as f:
            async for chunk in self._read_chunks(f):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()
    
    async def _read_chunks(self, file_obj, chunk_size: int = 8192):
        """异步读取文件块"""
        while True:
            chunk = await file_obj.read(chunk_size)
            if not chunk:
                break
            yield chunk
    
    async def save_uploaded_file(self, file_content: bytes, filename: str) -> Path:
        """保存上传的文件"""
        # 生成唯一文件名
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        file_ext = Path(filename).suffix
        unique_filename = f"{timestamp}_{filename}"
        file_path = self.upload_dir / unique_filename
        
        # 异步保存文件
        async with aiofiles.open(file_path, 'wb') as f:
            await f.write(file_content)
        
        logger.info(f"文件已保存: {file_path}")
        return file_path
    
    async def parse_pdf_with_marker(self, file_path: Path) -> Dict[str, Any]:
        """使用Marker解析PDF文档"""
        try:
            await self.initialize_models()
            
            # 在线程池中执行Marker解析，避免阻塞
            loop = asyncio.get_event_loop()
            result = await loop.run_in_executor(
                None, 
                self._parse_pdf_sync, 
                str(file_path)
            )
            
            return {
                "success": True,
                "content": result["markdown"],
                "metadata": {
                    "total_pages": result.get("pages", 0),
                    "processing_time": result.get("processing_time", 0),
                    "model_info": {
                        "detection_model": marker_settings.DETECTION_MODEL_NAME,
                        "recognition_model": marker_settings.RECOGNITION_MODEL_NAME,
                        "layout_model": marker_settings.LAYOUT_MODEL_NAME,
                    }
                }
            }
            
        except Exception as e:
            logger.error(f"Marker PDF解析失败: {e}")
            return {
                "success": False,
                "error": str(e),
                "content": "",
                "metadata": {}
            }
    
    def _parse_pdf_sync(self, file_path: str) -> Dict[str, Any]:
        """同步执行Marker PDF解析"""
        start_time = datetime.now()
        
        # 使用Marker转换PDF
        full_text, images, out_meta = convert_single_pdf(
            file_path, 
            self.models,
            max_pages=None,
            langs=["zh", "en"],  # 支持中英文
            batch_multiplier=2
        )
        
        processing_time = (datetime.now() - start_time).total_seconds()
        
        return {
            "markdown": full_text,
            "images": images,
            "metadata": out_meta,
            "pages": out_meta.get("page_count", 0),
            "processing_time": processing_time
        }
    
    async def parse_text_document(self, file_path: Path, file_ext: str) -> Dict[str, Any]:
        """解析文本类文档 (MD, TXT)"""
        try:
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
                content = await f.read()
            
            return {
                "success": True,
                "content": content,
                "metadata": {
                    "total_chars": len(content),
                    "total_lines": content.count('\n') + 1,
                    "encoding": "utf-8"
                }
            }
            
        except UnicodeDecodeError:
            # 尝试其他编码
            try:
                async with aiofiles.open(file_path, 'r', encoding='gbk') as f:
                    content = await f.read()
                
                return {
                    "success": True,
                    "content": content,
                    "metadata": {
                        "total_chars": len(content),
                        "total_lines": content.count('\n') + 1,
                        "encoding": "gbk"
                    }
                }
            except Exception as e:
                logger.error(f"文本文档解析失败: {e}")
                return {
                    "success": False,
                    "error": f"编码错误: {e}",
                    "content": "",
                    "metadata": {}
                }
        except Exception as e:
            logger.error(f"文本文档解析失败: {e}")
            return {
                "success": False,
                "error": str(e),
                "content": "",
                "metadata": {}
            }
    
    async def parse_office_document(self, file_path: Path, file_ext: str) -> Dict[str, Any]:
        """解析Office文档 (DOCX, PPTX, XLSX)"""
        try:
            if file_ext == 'docx':
                return await self._parse_docx(file_path)
            elif file_ext == 'pptx':
                return await self._parse_pptx(file_path)
            elif file_ext == 'xlsx':
                return await self._parse_xlsx(file_path)
            else:
                return {
                    "success": False,
                    "error": f"不支持的Office格式: {file_ext}",
                    "content": "",
                    "metadata": {}
                }
                
        except Exception as e:
            logger.error(f"Office文档解析失败: {e}")
            return {
                "success": False,
                "error": str(e),
                "content": "",
                "metadata": {}
            }
    
    async def _parse_docx(self, file_path: Path) -> Dict[str, Any]:
        """解析DOCX文档"""
        from docx import Document
        
        loop = asyncio.get_event_loop()
        doc = await loop.run_in_executor(None, Document, str(file_path))
        
        # 提取文本内容
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text.strip())
        
        content = '\n\n'.join(paragraphs)
        
        return {
            "success": True,
            "content": content,
            "metadata": {
                "total_paragraphs": len(paragraphs),
                "total_chars": len(content),
                "document_type": "docx"
            }
        }
    
    async def _parse_pptx(self, file_path: Path) -> Dict[str, Any]:
        """解析PPTX文档"""
        from pptx import Presentation
        
        loop = asyncio.get_event_loop()
        prs = await loop.run_in_executor(None, Presentation, str(file_path))
        
        # 提取幻灯片内容
        slides_content = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                slides_content.append(f"## 幻灯片 {i+1}\n\n" + '\n\n'.join(slide_text))
        
        content = '\n\n---\n\n'.join(slides_content)
        
        return {
            "success": True,
            "content": content,
            "metadata": {
                "total_slides": len(prs.slides),
                "total_chars": len(content),
                "document_type": "pptx"
            }
        }
    
    async def _parse_xlsx(self, file_path: Path) -> Dict[str, Any]:
        """解析XLSX文档"""
        import pandas as pd
        
        loop = asyncio.get_event_loop()
        
        # 读取所有工作表
        excel_file = await loop.run_in_executor(None, pd.ExcelFile, str(file_path))
        
        sheets_content = []
        for sheet_name in excel_file.sheet_names:
            df = await loop.run_in_executor(None, pd.read_excel, str(file_path), sheet_name)
            
            # 转换为Markdown表格格式
            if not df.empty:
                markdown_table = df.to_markdown(index=False)
                sheets_content.append(f"## 工作表: {sheet_name}\n\n{markdown_table}")
        
        content = '\n\n---\n\n'.join(sheets_content)
        
        return {
            "success": True,
            "content": content,
            "metadata": {
                "total_sheets": len(excel_file.sheet_names),
                "total_chars": len(content),
                "document_type": "xlsx"
            }
        }
    
    async def parse_document(self, file_path: Path, file_ext: str) -> Dict[str, Any]:
        """统一文档解析接口"""
        logger.info(f"开始解析文档: {file_path} (格式: {file_ext})")
        
        if file_ext == 'pdf':
            result = await self.parse_pdf_with_marker(file_path)
        elif file_ext in ['md', 'txt']:
            result = await self.parse_text_document(file_path, file_ext)
        elif file_ext in ['docx', 'pptx', 'xlsx']:
            result = await self.parse_office_document(file_path, file_ext)
        else:
            result = {
                "success": False,
                "error": f"不支持的文件格式: {file_ext}",
                "content": "",
                "metadata": {}
            }
        
        logger.info(f"文档解析完成: {file_path}, 成功: {result['success']}")
        return result
    
    async def cleanup_temp_file(self, file_path: Path):
        """清理临时文件"""
        try:
            if file_path.exists():
                await asyncio.get_event_loop().run_in_executor(None, file_path.unlink)
                logger.info(f"临时文件已删除: {file_path}")
        except Exception as e:
            logger.warning(f"删除临时文件失败: {e}")


# 全局文档解析服务实例
marker_service = MarkerDocumentService()
