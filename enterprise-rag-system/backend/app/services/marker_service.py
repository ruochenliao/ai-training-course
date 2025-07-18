"""
Marker文档解析服务 - 基于官方源码优化实现
支持高质量PDF解析和多种文档格式处理
"""

import asyncio
import os
import tempfile
import time
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Dict, Any, List

import fitz  # PyMuPDF
import pandas as pd
from docx import Document
from loguru import logger
from pptx import Presentation

# Marker imports - 基于官方源码
try:
    # 尝试导入marker-pdf包
    import marker
    from marker.convert import convert_single_pdf
    from marker.models import load_all_models
    from marker.settings import settings as marker_settings
    from marker.output import markdown_exists, save_markdown
    from marker.pdf.utils import find_filetype
    from marker.pdf.extract_text import get_length_of_text
    MARKER_AVAILABLE = True
    logger.info("Marker库加载成功")
except ImportError as e:
    logger.warning(f"Marker库不可用，将使用基础PDF解析: {e}")
    # 尝试安装marker-pdf
    try:
        import subprocess
        import sys
        logger.info("尝试安装marker-pdf...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "marker-pdf"])

        # 重新导入
        import marker
        from marker.convert import convert_single_pdf
        from marker.models import load_all_models
        from marker.settings import settings as marker_settings
        from marker.output import markdown_exists, save_markdown
        from marker.pdf.utils import find_filetype
        from marker.pdf.extract_text import get_length_of_text
        MARKER_AVAILABLE = True
        logger.info("Marker库安装并加载成功")
    except Exception as install_error:
        logger.warning(f"Marker库安装失败: {install_error}")
        MARKER_AVAILABLE = False

from app.core import settings
from app.core import DocumentProcessingException


class MarkerService:
    """Marker文档解析服务类 - 基于官方源码优化"""

    def __init__(self):
        """初始化Marker服务"""
        self.executor = ThreadPoolExecutor(max_workers=settings.MARKER_BATCH_MULTIPLIER or 2)
        self.supported_formats = {
            '.pdf': self._parse_pdf_with_marker,
            '.docx': self._parse_docx,
            '.doc': self._parse_docx,
            '.pptx': self._parse_pptx,
            '.ppt': self._parse_pptx,
            '.txt': self._parse_txt,
            '.md': self._parse_txt,
            '.html': self._parse_html,
            '.htm': self._parse_html,
            '.csv': self._parse_csv,
            '.xlsx': self._parse_excel,
            '.xls': self._parse_excel,
            '.json': self._parse_json
        }

        # Marker配置
        self.marker_enabled = settings.MARKER_ENABLED and MARKER_AVAILABLE
        self.marker_max_pages = settings.MARKER_MAX_PAGES
        self.marker_languages = settings.MARKER_LANGUAGES
        self.marker_batch_multiplier = settings.MARKER_BATCH_MULTIPLIER

        # 初始化Marker模型
        self.marker_models = None
        self._model_loading = False

        if self.marker_enabled:
            asyncio.create_task(self._initialize_marker_models())
        else:
            self.supported_formats['.pdf'] = self._parse_pdf_basic
            logger.info("Marker已禁用，使用基础PDF解析")

        logger.info("Marker文档解析服务初始化完成")

    async def _initialize_marker_models(self):
        """异步初始化Marker模型"""
        if self._model_loading:
            return

        self._model_loading = True
        try:
            logger.info("开始加载Marker模型...")
            loop = asyncio.get_event_loop()
            self.marker_models = await loop.run_in_executor(
                self.executor,
                self._load_marker_models_sync
            )
            logger.info("Marker模型加载成功")
        except Exception as e:
            logger.warning(f"Marker模型加载失败，将使用基础PDF解析: {e}")
            self.supported_formats['.pdf'] = self._parse_pdf_basic
            self.marker_enabled = False
        finally:
            self._model_loading = False

    def _load_marker_models_sync(self):
        """同步加载Marker模型"""
        try:
            # 设置Marker配置
            if hasattr(marker_settings, 'TORCH_DEVICE_MODEL'):
                marker_settings.TORCH_DEVICE_MODEL = "cpu"  # 可根据需要调整

            models = load_all_models()
            return models
        except Exception as e:
            logger.error(f"同步加载Marker模型失败: {e}")
            raise
    
    async def parse_document(
        self,
        file_path: str,
        file_name: str,
        extract_images: bool = True,
        extract_tables: bool = True,
        force_ocr: bool = False
    ) -> Dict[str, Any]:
        """
        解析文档 - 支持批量处理和进度跟踪

        Args:
            file_path: 文件路径
            file_name: 文件名
            extract_images: 是否提取图片
            extract_tables: 是否提取表格
            force_ocr: 是否强制使用OCR

        Returns:
            解析结果字典
        """
        start_time = time.time()

        try:
            # 验证文件
            if not os.path.exists(file_path):
                raise DocumentProcessingException(f"文件不存在: {file_path}")

            file_size = os.path.getsize(file_path)
            if file_size > settings.MAX_FILE_SIZE:
                raise DocumentProcessingException(f"文件过大: {file_size} bytes")

            # 获取文件扩展名
            file_ext = Path(file_name).suffix.lower()

            if file_ext not in self.supported_formats:
                raise DocumentProcessingException(f"不支持的文件格式: {file_ext}")

            logger.info(f"开始解析文档: {file_name} ({file_size} bytes)")

            # 在线程池中执行解析
            loop = asyncio.get_event_loop()
            result = await loop.run_in_executor(
                self.executor,
                self._parse_document_sync,
                file_path,
                file_name,
                file_ext,
                extract_images,
                extract_tables,
                force_ocr
            )

            # 添加处理时间
            processing_time = time.time() - start_time
            result['processing_time'] = processing_time
            result['file_size'] = file_size

            logger.info(f"文档解析完成: {file_name} (耗时: {processing_time:.2f}s)")
            return result

        except Exception as e:
            logger.error(f"文档解析失败 {file_name}: {e}")
            raise DocumentProcessingException(f"文档解析失败: {e}")

    async def parse_documents_batch(
        self,
        file_paths: List[str],
        extract_images: bool = True,
        extract_tables: bool = True,
        max_concurrent: int = 3
    ) -> List[Dict[str, Any]]:
        """
        批量解析文档

        Args:
            file_paths: 文件路径列表
            extract_images: 是否提取图片
            extract_tables: 是否提取表格
            max_concurrent: 最大并发数

        Returns:
            解析结果列表
        """
        semaphore = asyncio.Semaphore(max_concurrent)

        async def parse_single(file_path: str):
            async with semaphore:
                file_name = os.path.basename(file_path)
                try:
                    return await self.parse_document(
                        file_path, file_name, extract_images, extract_tables
                    )
                except Exception as e:
                    logger.error(f"批量解析失败 {file_name}: {e}")
                    return {
                        'file_name': file_name,
                        'file_path': file_path,
                        'error': str(e),
                        'success': False
                    }

        logger.info(f"开始批量解析 {len(file_paths)} 个文档")
        results = await asyncio.gather(*[parse_single(fp) for fp in file_paths])

        success_count = len([r for r in results if r.get('success', True)])
        logger.info(f"批量解析完成: {success_count}/{len(file_paths)} 成功")

        return results
    
    def _parse_document_sync(
        self,
        file_path: str,
        file_name: str,
        file_ext: str,
        extract_images: bool,
        extract_tables: bool,
        force_ocr: bool = False
    ) -> Dict[str, Any]:
        """同步解析文档"""
        parser_func = self.supported_formats[file_ext]

        result = {
            'file_name': file_name,
            'file_path': file_path,
            'file_type': file_ext,
            'content': '',
            'metadata': {},
            'images': [],
            'tables': [],
            'structure': {},
            'page_count': 0,
            'word_count': 0,
            'success': True,
            'error': None
        }

        try:
            # 调用对应的解析函数
            parsed_data = parser_func(file_path, extract_images, extract_tables, force_ocr)
            result.update(parsed_data)

            # 计算字数和字符数
            if result['content']:
                content = result['content']
                result['word_count'] = len(content.split())
                result['char_count'] = len(content)
                result['line_count'] = len(content.split('\n'))

            # 添加质量评估
            result['quality_score'] = self._assess_content_quality(result)

        except Exception as e:
            result['success'] = False
            result['error'] = str(e)
            logger.error(f"文档解析失败 {file_name}: {e}")

        return result

    def _assess_content_quality(self, result: Dict[str, Any]) -> float:
        """评估内容质量"""
        score = 0.0

        # 基础分数
        if result.get('content'):
            score += 0.3

        # 结构化程度
        if result.get('tables'):
            score += 0.2

        if result.get('images'):
            score += 0.1

        # 元数据完整性
        metadata = result.get('metadata', {})
        if metadata.get('title'):
            score += 0.1
        if metadata.get('author'):
            score += 0.1

        # 内容长度合理性
        word_count = result.get('word_count', 0)
        if 100 <= word_count <= 50000:
            score += 0.2
        elif word_count > 0:
            score += 0.1

        return min(score, 1.0)

    def _parse_pdf_with_marker(self, file_path: str, extract_images: bool, extract_tables: bool, force_ocr: bool = False) -> Dict[str, Any]:
        """使用Marker解析PDF文件（基于官方源码优化）"""
        if not self.marker_enabled or not self.marker_models:
            logger.info("Marker不可用，使用基础PDF解析")
            return self._parse_pdf_basic(file_path, extract_images, extract_tables, force_ocr)

        try:
            logger.info(f"使用Marker解析PDF: {file_path}")

            # 检查文件类型
            if MARKER_AVAILABLE:
                try:
                    filetype = find_filetype(file_path)
                    if filetype != "pdf":
                        logger.warning(f"文件类型不匹配: {filetype}")
                        return self._parse_pdf_basic(file_path, extract_images, extract_tables, force_ocr)
                except:
                    pass

            # 检查文档长度（避免处理过长文档）
            try:
                text_length = get_length_of_text(file_path)
                if text_length > 500000:  # 50万字符限制
                    logger.warning(f"文档过长 ({text_length} 字符)，使用基础解析")
                    return self._parse_pdf_basic(file_path, extract_images, extract_tables, force_ocr)
            except:
                pass

            # 使用Marker进行高质量PDF解析
            start_time = time.time()

            full_text, images, out_meta = convert_single_pdf(
                file_path,
                self.marker_models,
                max_pages=self.marker_max_pages,
                langs=self.marker_languages,
                batch_multiplier=self.marker_batch_multiplier,
                start_page=None
            )

            processing_time = time.time() - start_time
            logger.info(f"Marker解析完成，耗时: {processing_time:.2f}s")

            # 处理图片
            processed_images = []
            if extract_images and images:
                for img_key, img_data in images.items():
                    try:
                        # 保存图片到临时文件
                        temp_dir = tempfile.gettempdir()
                        img_path = os.path.join(temp_dir, f"marker_img_{img_key}.png")

                        with open(img_path, 'wb') as f:
                            f.write(img_data)

                        processed_images.append({
                            'key': img_key,
                            'path': img_path,
                            'size': len(img_data),
                            'type': 'marker_extracted',
                            'format': 'png'
                        })
                    except Exception as e:
                        logger.warning(f"处理图片失败 {img_key}: {e}")

            # 从元数据中提取表格信息
            tables = []
            if extract_tables and out_meta:
                # 尝试从不同字段提取表格
                table_sources = ['table_blocks', 'tables', 'table_data']
                for source in table_sources:
                    if source in out_meta and out_meta[source]:
                        for i, table_data in enumerate(out_meta[source]):
                            try:
                                table_info = {
                                    'index': i,
                                    'source': source,
                                    'type': 'marker_table'
                                }

                                if isinstance(table_data, dict):
                                    table_info.update({
                                        'bbox': table_data.get('bbox', []),
                                        'content': table_data.get('text', ''),
                                        'confidence': table_data.get('confidence', 0.0)
                                    })
                                else:
                                    table_info['content'] = str(table_data)

                                tables.append(table_info)
                            except Exception as e:
                                logger.warning(f"处理表格失败 {i}: {e}")
                        break

            # 构建结果
            return {
                'content': full_text,
                'metadata': {
                    'marker_version': True,
                    'marker_processing_time': processing_time,
                    'languages': out_meta.get('languages', self.marker_languages),
                    'page_count': out_meta.get('page_count', 0),
                    'text_length': len(full_text),
                    'model_info': {
                        'detection_model': out_meta.get('detection_model', ''),
                        'recognition_model': out_meta.get('recognition_model', ''),
                        'layout_model': out_meta.get('layout_model', '')
                    }
                },
                'images': processed_images,
                'tables': tables,
                'page_count': out_meta.get('page_count', 0),
                'structure': {
                    'type': 'pdf_marker',
                    'blocks': len(out_meta.get('blocks', [])),
                    'table_blocks': len(tables),
                    'image_blocks': len(processed_images),
                    'total_elements': len(out_meta.get('blocks', []))
                }
            }

        except Exception as e:
            logger.warning(f"Marker PDF解析失败，回退到基础解析: {e}")
            return self._parse_pdf_basic(file_path, extract_images, extract_tables, force_ocr)

    def _parse_pdf_basic(self, file_path: str, extract_images: bool, extract_tables: bool, force_ocr: bool = False) -> Dict[str, Any]:
        """基础PDF解析（使用PyMuPDF）"""
        try:
            doc = fitz.open(file_path)
            content_parts = []
            images = []
            tables = []
            metadata = {}
            
            # 获取文档元数据
            metadata.update({
                'title': doc.metadata.get('title', ''),
                'author': doc.metadata.get('author', ''),
                'subject': doc.metadata.get('subject', ''),
                'creator': doc.metadata.get('creator', ''),
                'producer': doc.metadata.get('producer', ''),
                'creation_date': doc.metadata.get('creationDate', ''),
                'modification_date': doc.metadata.get('modDate', '')
            })
            
            # 逐页处理
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # 提取文本
                text = page.get_text()
                if text.strip():
                    content_parts.append(f"## 第{page_num + 1}页\n\n{text}\n")
                
                # 提取图片
                if extract_images:
                    image_list = page.get_images()
                    for img_index, img in enumerate(image_list):
                        try:
                            xref = img[0]
                            pix = fitz.Pixmap(doc, xref)
                            if pix.n - pix.alpha < 4:  # 确保是RGB或灰度图
                                img_data = pix.tobytes("png")
                                images.append({
                                    'page': page_num + 1,
                                    'index': img_index,
                                    'data': img_data,
                                    'width': pix.width,
                                    'height': pix.height
                                })
                            pix = None
                        except Exception as e:
                            logger.warning(f"提取图片失败 page {page_num + 1}, img {img_index}: {e}")
                
                # 提取表格（简单实现）
                if extract_tables:
                    tables_on_page = page.find_tables()
                    for table_index, table in enumerate(tables_on_page):
                        try:
                            table_data = table.extract()
                            tables.append({
                                'page': page_num + 1,
                                'index': table_index,
                                'data': table_data,
                                'bbox': table.bbox
                            })
                        except Exception as e:
                            logger.warning(f"提取表格失败 page {page_num + 1}, table {table_index}: {e}")
            
            doc.close()
            
            return {
                'content': '\n'.join(content_parts),
                'metadata': metadata,
                'images': images,
                'tables': tables,
                'page_count': len(doc),
                'structure': {'type': 'pdf', 'pages': len(doc)}
            }
            
        except Exception as e:
            logger.error(f"PDF解析失败: {e}")
            raise DocumentProcessingException(f"PDF解析失败: {e}")
    
    def _parse_docx(self, file_path: str, extract_images: bool, extract_tables: bool, force_ocr: bool = False) -> Dict[str, Any]:
        """解析DOCX文件"""
        try:
            doc = Document(file_path)
            content_parts = []
            images = []
            tables = []
            
            # 获取文档属性
            metadata = {
                'title': doc.core_properties.title or '',
                'author': doc.core_properties.author or '',
                'subject': doc.core_properties.subject or '',
                'created': str(doc.core_properties.created) if doc.core_properties.created else '',
                'modified': str(doc.core_properties.modified) if doc.core_properties.modified else ''
            }
            
            # 提取段落文本
            for para in doc.paragraphs:
                if para.text.strip():
                    # 根据样式添加标记
                    if para.style.name.startswith('Heading'):
                        level = para.style.name.replace('Heading ', '')
                        content_parts.append(f"{'#' * int(level)} {para.text}\n")
                    else:
                        content_parts.append(f"{para.text}\n")
            
            # 提取表格
            if extract_tables:
                for table_index, table in enumerate(doc.tables):
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    
                    tables.append({
                        'index': table_index,
                        'data': table_data,
                        'rows': len(table.rows),
                        'cols': len(table.columns) if table.rows else 0
                    })
            
            # 提取图片（简化实现）
            if extract_images:
                # 这里可以扩展图片提取逻辑
                pass
            
            return {
                'content': '\n'.join(content_parts),
                'metadata': metadata,
                'images': images,
                'tables': tables,
                'page_count': 1,  # Word文档没有固定页数概念
                'structure': {'type': 'docx', 'paragraphs': len(doc.paragraphs), 'tables': len(doc.tables)}
            }
            
        except Exception as e:
            logger.error(f"DOCX解析失败: {e}")
            raise DocumentProcessingException(f"DOCX解析失败: {e}")
    
    def _parse_pptx(self, file_path: str, extract_images: bool, extract_tables: bool, force_ocr: bool = False) -> Dict[str, Any]:
        """解析PPTX文件"""
        try:
            prs = Presentation(file_path)
            content_parts = []
            images = []
            tables = []
            
            # 获取演示文稿属性
            metadata = {
                'title': prs.core_properties.title or '',
                'author': prs.core_properties.author or '',
                'subject': prs.core_properties.subject or '',
                'created': str(prs.core_properties.created) if prs.core_properties.created else '',
                'modified': str(prs.core_properties.modified) if prs.core_properties.modified else ''
            }
            
            # 逐幻灯片处理
            for slide_index, slide in enumerate(prs.slides):
                slide_content = [f"## 幻灯片 {slide_index + 1}\n"]
                
                # 提取文本
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)
                    
                    # 提取表格
                    if extract_tables and shape.shape_type == 19:  # 表格类型
                        try:
                            table = shape.table
                            table_data = []
                            for row in table.rows:
                                row_data = [cell.text.strip() for cell in row.cells]
                                table_data.append(row_data)
                            
                            tables.append({
                                'slide': slide_index + 1,
                                'data': table_data,
                                'rows': len(table.rows),
                                'cols': len(table.rows[0].cells) if table.rows else 0
                            })
                        except Exception as e:
                            logger.warning(f"提取PPT表格失败 slide {slide_index + 1}: {e}")
                
                content_parts.append('\n'.join(slide_content) + '\n')
            
            return {
                'content': '\n'.join(content_parts),
                'metadata': metadata,
                'images': images,
                'tables': tables,
                'page_count': len(prs.slides),
                'structure': {'type': 'pptx', 'slides': len(prs.slides)}
            }
            
        except Exception as e:
            logger.error(f"PPTX解析失败: {e}")
            raise DocumentProcessingException(f"PPTX解析失败: {e}")
    
    def _parse_txt(self, file_path: str, extract_images: bool, extract_tables: bool, force_ocr: bool = False) -> Dict[str, Any]:
        """解析文本文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            return {
                'content': content,
                'metadata': {'encoding': 'utf-8'},
                'images': [],
                'tables': [],
                'page_count': 1,
                'structure': {'type': 'text', 'lines': len(content.split('\n'))}
            }
            
        except UnicodeDecodeError:
            # 尝试其他编码
            for encoding in ['gbk', 'gb2312', 'latin1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    
                    return {
                        'content': content,
                        'metadata': {'encoding': encoding},
                        'images': [],
                        'tables': [],
                        'page_count': 1,
                        'structure': {'type': 'text', 'lines': len(content.split('\n'))}
                    }
                except UnicodeDecodeError:
                    continue
            
            raise DocumentProcessingException("无法识别文本文件编码")
            
        except Exception as e:
            logger.error(f"文本文件解析失败: {e}")
            raise DocumentProcessingException(f"文本文件解析失败: {e}")
    
    def _parse_html(self, file_path: str, extract_images: bool, extract_tables: bool) -> Dict[str, Any]:
        """解析HTML文件"""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # 提取文本内容
            text_content = soup.get_text(separator='\n', strip=True)
            
            # 提取元数据
            metadata = {}
            title_tag = soup.find('title')
            if title_tag:
                metadata['title'] = title_tag.get_text()
            
            meta_tags = soup.find_all('meta')
            for meta in meta_tags:
                if meta.get('name'):
                    metadata[meta.get('name')] = meta.get('content', '')
            
            # 提取表格
            tables = []
            if extract_tables:
                table_tags = soup.find_all('table')
                for table_index, table in enumerate(table_tags):
                    table_data = []
                    rows = table.find_all('tr')
                    for row in rows:
                        cells = row.find_all(['td', 'th'])
                        row_data = [cell.get_text(strip=True) for cell in cells]
                        if row_data:
                            table_data.append(row_data)
                    
                    if table_data:
                        tables.append({
                            'index': table_index,
                            'data': table_data,
                            'rows': len(table_data),
                            'cols': len(table_data[0]) if table_data else 0
                        })
            
            return {
                'content': text_content,
                'metadata': metadata,
                'images': [],
                'tables': tables,
                'page_count': 1,
                'structure': {'type': 'html', 'tags': len(soup.find_all())}
            }
            
        except Exception as e:
            logger.error(f"HTML解析失败: {e}")
            raise DocumentProcessingException(f"HTML解析失败: {e}")
    
    def _parse_csv(self, file_path: str, extract_images: bool, extract_tables: bool) -> Dict[str, Any]:
        """解析CSV文件"""
        try:
            df = pd.read_csv(file_path)
            
            # 转换为文本格式
            content = df.to_string(index=False)
            
            # 表格数据
            table_data = [df.columns.tolist()] + df.values.tolist()
            
            return {
                'content': content,
                'metadata': {'rows': len(df), 'columns': len(df.columns)},
                'images': [],
                'tables': [{
                    'index': 0,
                    'data': table_data,
                    'rows': len(df) + 1,  # +1 for header
                    'cols': len(df.columns)
                }],
                'page_count': 1,
                'structure': {'type': 'csv', 'rows': len(df), 'columns': len(df.columns)}
            }
            
        except Exception as e:
            logger.error(f"CSV解析失败: {e}")
            raise DocumentProcessingException(f"CSV解析失败: {e}")
    
    def _parse_excel(self, file_path: str, extract_images: bool, extract_tables: bool) -> Dict[str, Any]:
        """解析Excel文件"""
        try:
            # 读取所有工作表
            excel_file = pd.ExcelFile(file_path)
            content_parts = []
            tables = []
            
            for sheet_index, sheet_name in enumerate(excel_file.sheet_names):
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # 添加工作表标题
                content_parts.append(f"## 工作表: {sheet_name}\n")
                content_parts.append(df.to_string(index=False))
                content_parts.append('\n')
                
                # 表格数据
                if not df.empty:
                    table_data = [df.columns.tolist()] + df.values.tolist()
                    tables.append({
                        'sheet': sheet_name,
                        'index': sheet_index,
                        'data': table_data,
                        'rows': len(df) + 1,
                        'cols': len(df.columns)
                    })
            
            return {
                'content': '\n'.join(content_parts),
                'metadata': {'sheets': len(excel_file.sheet_names)},
                'images': [],
                'tables': tables,
                'page_count': len(excel_file.sheet_names),
                'structure': {'type': 'excel', 'sheets': len(excel_file.sheet_names)}
            }
            
        except Exception as e:
            logger.error(f"Excel解析失败: {e}")
            raise DocumentProcessingException(f"Excel解析失败: {e}")
    
    def _parse_json(self, file_path: str, extract_images: bool, extract_tables: bool) -> Dict[str, Any]:
        """解析JSON文件"""
        try:
            import json
            
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # 格式化JSON为可读文本
            content = json.dumps(data, ensure_ascii=False, indent=2)
            
            return {
                'content': content,
                'metadata': {'type': type(data).__name__},
                'images': [],
                'tables': [],
                'page_count': 1,
                'structure': {'type': 'json', 'size': len(str(data))}
            }
            
        except Exception as e:
            logger.error(f"JSON解析失败: {e}")
            raise DocumentProcessingException(f"JSON解析失败: {e}")


# 全局Marker服务实例
marker_service = MarkerService()
