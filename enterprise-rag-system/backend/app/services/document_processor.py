"""
文档处理服务
"""

import hashlib
import mimetypes
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any

import aiofiles
from loguru import logger

from app.core import DocumentProcessingException
from app.core import settings
from app.models import Document, DocumentChunk
from app.schemas import ProcessingResult, DocumentMetadata


class DocumentProcessor:
    """文档处理器基类"""
    
    def __init__(self):
        self.supported_formats = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.pptx': self._process_pptx,
            '.txt': self._process_text,
            '.md': self._process_markdown,
            '.html': self._process_html,
            '.csv': self._process_csv,
            '.xlsx': self._process_excel,
            '.json': self._process_json
        }
    
    async def process_document(
        self, 
        file_path: str, 
        document_id: int,
        knowledge_base_id: int
    ) -> ProcessingResult:
        """处理文档"""
        try:
            # 获取文件信息
            file_info = await self._get_file_info(file_path)
            file_ext = file_info['extension'].lower()
            
            # 检查文件格式支持
            if file_ext not in self.supported_formats:
                raise DocumentProcessingException(f"不支持的文件格式: {file_ext}")
            
            # 更新文档状态为处理中
            document = await Document.get(id=document_id)
            await document.start_processing()
            
            # 选择处理器
            processor = self.supported_formats[file_ext]
            
            # 处理文档
            result = await processor(file_path, file_info)
            
            # 智能分块
            chunks = await self._chunk_content(
                result['content'], 
                document_id,
                knowledge_base_id,
                result.get('metadata', {})
            )
            
            # 保存处理结果
            await self._save_processing_result(document, result, chunks)
            
            # 更新文档状态为完成
            await document.complete_processing(len(chunks))
            
            logger.info(f"文档处理完成: {file_path}, 生成 {len(chunks)} 个分块")
            
            return ProcessingResult(
                success=True,
                content=result['content'],
                metadata=result.get('metadata', {}),
                chunks=len(chunks),
                tables=result.get('tables', []),
                images=result.get('images', [])
            )
            
        except Exception as e:
            # 更新文档状态为失败
            try:
                document = await Document.get(id=document_id)
                await document.fail_processing(str(e))
            except:
                pass
            
            logger.error(f"文档处理失败: {file_path}, 错误: {e}")
            raise DocumentProcessingException(f"文档处理失败: {e}")
    
    async def _get_file_info(self, file_path: str) -> Dict[str, Any]:
        """获取文件信息"""
        try:
            path = Path(file_path)
            stat = path.stat()
            
            # 计算文件哈希
            file_hash = await self._calculate_file_hash(file_path)
            
            # 获取MIME类型
            mime_type, _ = mimetypes.guess_type(file_path)
            
            return {
                'name': path.name,
                'extension': path.suffix,
                'size': stat.st_size,
                'modified_time': datetime.fromtimestamp(stat.st_mtime),
                'hash': file_hash,
                'mime_type': mime_type
            }
        except Exception as e:
            raise DocumentProcessingException(f"获取文件信息失败: {e}")
    
    async def _calculate_file_hash(self, file_path: str) -> str:
        """计算文件哈希值"""
        hash_md5 = hashlib.md5()
        async with aiofiles.open(file_path, 'rb') as f:
            async for chunk in f:
                hash_md5.update(chunk)
        return hash_md5.hexdigest()
    
    async def _process_pdf(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """处理PDF文件"""
        try:
            # 这里应该集成Marker库
            # 由于Marker可能需要特殊环境，我们先用简单的PDF处理
            import PyPDF2
            
            content = ""
            tables = []
            images = []
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    content += f"\n\n--- 第 {page_num + 1} 页 ---\n\n"
                    content += page_text
            
            metadata = {
                'pages': len(pdf_reader.pages),
                'file_info': file_info,
                'processing_method': 'PyPDF2'
            }
            
            return {
                'content': content.strip(),
                'metadata': metadata,
                'tables': tables,
                'images': images
            }
            
        except Exception as e:
            raise DocumentProcessingException(f"PDF处理失败: {e}")
    
    async def _process_docx(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """处理Word文档"""
        try:
            from docx import Document as DocxDocument
            
            doc = DocxDocument(file_path)
            content = ""
            tables = []
            
            # 提取段落
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content += paragraph.text + "\n\n"
            
            # 提取表格
            for table_idx, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data.append(row_data)
                
                tables.append({
                    'index': table_idx,
                    'data': table_data,
                    'rows': len(table_data),
                    'cols': len(table_data[0]) if table_data else 0
                })
            
            metadata = {
                'paragraphs': len(doc.paragraphs),
                'tables': len(tables),
                'file_info': file_info,
                'processing_method': 'python-docx'
            }
            
            return {
                'content': content.strip(),
                'metadata': metadata,
                'tables': tables,
                'images': []
            }
            
        except Exception as e:
            raise DocumentProcessingException(f"Word文档处理失败: {e}")
    
    async def _process_pptx(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """处理PowerPoint文档"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            content = ""
            
            for slide_idx, slide in enumerate(prs.slides):
                content += f"\n\n--- 幻灯片 {slide_idx + 1} ---\n\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content += shape.text + "\n"
            
            metadata = {
                'slides': len(prs.slides),
                'file_info': file_info,
                'processing_method': 'python-pptx'
            }
            
            return {
                'content': content.strip(),
                'metadata': metadata,
                'tables': [],
                'images': []
            }
            
        except Exception as e:
            raise DocumentProcessingException(f"PowerPoint文档处理失败: {e}")
    
    async def _process_text(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """处理纯文本文件"""
        try:
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
                content = await f.read()
            
            metadata = {
                'encoding': 'utf-8',
                'file_info': file_info,
                'processing_method': 'text'
            }
            
            return {
                'content': content,
                'metadata': metadata,
                'tables': [],
                'images': []
            }
            
        except UnicodeDecodeError:
            # 尝试其他编码
            try:
                async with aiofiles.open(file_path, 'r', encoding='gbk') as f:
                    content = await f.read()
                
                metadata = {
                    'encoding': 'gbk',
                    'file_info': file_info,
                    'processing_method': 'text'
                }
                
                return {
                    'content': content,
                    'metadata': metadata,
                    'tables': [],
                    'images': []
                }
            except Exception as e:
                raise DocumentProcessingException(f"文本文件编码识别失败: {e}")
        except Exception as e:
            raise DocumentProcessingException(f"文本文件处理失败: {e}")
    
    async def _process_markdown(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """处理Markdown文件"""
        try:
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
                content = await f.read()
            
            metadata = {
                'format': 'markdown',
                'file_info': file_info,
                'processing_method': 'markdown'
            }
            
            return {
                'content': content,
                'metadata': metadata,
                'tables': [],
                'images': []
            }
            
        except Exception as e:
            raise DocumentProcessingException(f"Markdown文件处理失败: {e}")
    
    async def _process_html(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """处理HTML文件"""
        try:
            from bs4 import BeautifulSoup
            
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
                html_content = await f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # 移除脚本和样式
            for script in soup(["script", "style"]):
                script.decompose()
            
            # 提取文本
            content = soup.get_text()
            
            # 清理空白
            lines = (line.strip() for line in content.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            content = '\n'.join(chunk for chunk in chunks if chunk)
            
            metadata = {
                'format': 'html',
                'title': soup.title.string if soup.title else None,
                'file_info': file_info,
                'processing_method': 'beautifulsoup'
            }
            
            return {
                'content': content,
                'metadata': metadata,
                'tables': [],
                'images': []
            }
            
        except Exception as e:
            raise DocumentProcessingException(f"HTML文件处理失败: {e}")
    
    async def _process_csv(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """处理CSV文件"""
        try:
            import pandas as pd
            
            df = pd.read_csv(file_path)
            
            # 转换为文本描述
            content = f"CSV文件包含 {len(df)} 行数据，{len(df.columns)} 个字段。\n\n"
            content += f"字段名称: {', '.join(df.columns)}\n\n"
            
            # 添加数据样本
            content += "数据样本:\n"
            content += df.head(10).to_string(index=False)
            
            # 添加统计信息
            content += "\n\n数据统计:\n"
            content += df.describe().to_string()
            
            metadata = {
                'rows': len(df),
                'columns': len(df.columns),
                'column_names': list(df.columns),
                'file_info': file_info,
                'processing_method': 'pandas'
            }
            
            return {
                'content': content,
                'metadata': metadata,
                'tables': [{'data': df.values.tolist(), 'headers': list(df.columns)}],
                'images': []
            }
            
        except Exception as e:
            raise DocumentProcessingException(f"CSV文件处理失败: {e}")
    
    async def _process_excel(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """处理Excel文件"""
        try:
            import pandas as pd
            
            # 读取所有工作表
            excel_file = pd.ExcelFile(file_path)
            content = f"Excel文件包含 {len(excel_file.sheet_names)} 个工作表。\n\n"
            
            tables = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                content += f"工作表: {sheet_name}\n"
                content += f"包含 {len(df)} 行数据，{len(df.columns)} 个字段。\n"
                content += f"字段: {', '.join(df.columns)}\n\n"
                
                # 添加数据样本
                if not df.empty:
                    content += df.head(5).to_string(index=False) + "\n\n"
                
                tables.append({
                    'sheet_name': sheet_name,
                    'data': df.values.tolist(),
                    'headers': list(df.columns),
                    'rows': len(df),
                    'cols': len(df.columns)
                })
            
            metadata = {
                'sheets': len(excel_file.sheet_names),
                'sheet_names': excel_file.sheet_names,
                'file_info': file_info,
                'processing_method': 'pandas'
            }
            
            return {
                'content': content,
                'metadata': metadata,
                'tables': tables,
                'images': []
            }
            
        except Exception as e:
            raise DocumentProcessingException(f"Excel文件处理失败: {e}")
    
    async def _process_json(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """处理JSON文件"""
        try:
            import json
            
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
                json_content = await f.read()
            
            data = json.loads(json_content)
            
            # 转换为可读文本
            content = f"JSON文件内容:\n\n"
            content += json.dumps(data, ensure_ascii=False, indent=2)
            
            metadata = {
                'format': 'json',
                'structure': type(data).__name__,
                'file_info': file_info,
                'processing_method': 'json'
            }
            
            return {
                'content': content,
                'metadata': metadata,
                'tables': [],
                'images': []
            }
            
        except Exception as e:
            raise DocumentProcessingException(f"JSON文件处理失败: {e}")
    
    async def _chunk_content(
        self, 
        content: str, 
        document_id: int,
        knowledge_base_id: int,
        metadata: Dict[str, Any]
    ) -> List[DocumentChunk]:
        """智能分块"""
        try:
            from app import DocumentChunker
            
            chunker = DocumentChunker()
            chunks = await chunker.chunk_text(
                content, 
                document_id,
                knowledge_base_id,
                metadata
            )
            
            return chunks
            
        except Exception as e:
            logger.error(f"文档分块失败: {e}")
            # 如果智能分块失败，使用简单分块
            return await self._simple_chunk(content, document_id, knowledge_base_id, metadata)
    
    async def _simple_chunk(
        self, 
        content: str, 
        document_id: int,
        knowledge_base_id: int,
        metadata: Dict[str, Any]
    ) -> List[DocumentChunk]:
        """简单分块"""
        chunk_size = settings.CHUNK_SIZE
        chunk_overlap = settings.CHUNK_OVERLAP
        
        chunks = []
        start = 0
        chunk_index = 0
        
        while start < len(content):
            end = start + chunk_size
            chunk_content = content[start:end]
            
            # 创建分块
            chunk = await DocumentChunk.create(
                document_id=document_id,
                chunk_index=chunk_index,
                content=chunk_content,
                content_hash=hashlib.md5(chunk_content.encode()).hexdigest(),
                start_position=start,
                end_position=end,
                char_count=len(chunk_content),
                word_count=len(chunk_content.split()),
                metadata=metadata
            )
            
            chunks.append(chunk)
            
            # 移动到下一个分块
            start = end - chunk_overlap
            chunk_index += 1
        
        return chunks
    
    async def _save_processing_result(
        self, 
        document: Document, 
        result: Dict[str, Any], 
        chunks: List[DocumentChunk]
    ):
        """保存处理结果"""
        try:
            # 更新文档信息
            document.extracted_text = result['content']
            document.extracted_tables = result.get('tables', [])
            document.extracted_images = result.get('images', [])
            document.chunk_count = len(chunks)
            
            # 更新元数据
            if result.get('metadata'):
                document.update_metadata(result['metadata'])
            
            await document.save()
            
        except Exception as e:
            logger.error(f"保存处理结果失败: {e}")
            raise DocumentProcessingException(f"保存处理结果失败: {e}")


# 全局文档处理器实例
document_processor = DocumentProcessor()
