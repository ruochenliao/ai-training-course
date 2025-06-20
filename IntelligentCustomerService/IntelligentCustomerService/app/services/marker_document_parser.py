"""
Marker文档解析服务 - 智能客服系统专用版
基于Marker框架的高质量文档解析，支持PDF、Word、Excel等多种格式
"""

import os
import asyncio
import tempfile
from typing import Dict, Any, Optional, List, Union, Tuple
from pathlib import Path
import json
import hashlib
from concurrent.futures import ThreadPoolExecutor
from dataclasses import dataclass

from loguru import logger
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd
import openpyxl
from PIL import Image
import io
import base64

from app.core.config import settings
from app.core.exceptions import DocumentProcessingException


@dataclass
class DocumentParseResult:
    """文档解析结果"""
    content: str
    metadata: Dict[str, Any]
    images: List[Dict[str, Any]]
    tables: List[Dict[str, Any]]
    structure: Dict[str, Any]
    page_count: int
    word_count: int
    file_hash: str
    processing_time: float


class MarkerDocumentParser:
    """
    Marker文档解析器
    支持多种文档格式的高质量解析
    """
    
    def __init__(self):
        """初始化解析器"""
        self.executor = ThreadPoolExecutor(max_workers=settings.DOCUMENT_PROCESSING_WORKERS)
        self.cache_dir = Path(settings.UPLOAD_DIR) / "document_cache"
        self.cache_dir.mkdir(exist_ok=True)
        
        # 支持的文档格式
        self.supported_formats = {
            '.pdf': self._parse_pdf_with_marker,
            '.docx': self._parse_docx,
            '.doc': self._parse_docx,
            '.pptx': self._parse_pptx,
            '.ppt': self._parse_pptx,
            '.txt': self._parse_text,
            '.md': self._parse_markdown,
            '.html': self._parse_html,
            '.htm': self._parse_html,
            '.csv': self._parse_csv,
            '.xlsx': self._parse_excel,
            '.xls': self._parse_excel,
            '.json': self._parse_json,
            '.xml': self._parse_xml,
            '.rtf': self._parse_rtf
        }
        
        # Marker模型配置
        self.marker_config = {
            'extract_images': True,
            'extract_tables': True,
            'ocr_all_pages': False,  # 仅在需要时使用OCR
            'languages': ['zh', 'en'],  # 支持中英文
            'batch_size': 1,
            'max_pages': 1000
        }
        
        logger.info("Marker文档解析器初始化完成")
    
    async def parse_document(
        self,
        file_path: str,
        file_name: str,
        extract_images: bool = True,
        extract_tables: bool = True,
        use_cache: bool = True
    ) -> DocumentParseResult:
        """
        解析文档
        
        Args:
            file_path: 文件路径
            file_name: 文件名
            extract_images: 是否提取图片
            extract_tables: 是否提取表格
            use_cache: 是否使用缓存
            
        Returns:
            DocumentParseResult: 解析结果
        """
        import time
        start_time = time.time()
        
        try:
            # 计算文件哈希
            file_hash = await self._calculate_file_hash(file_path)
            
            # 检查缓存
            if use_cache:
                cached_result = await self._get_cached_result(file_hash)
                if cached_result:
                    logger.info(f"使用缓存结果: {file_name}")
                    return cached_result
            
            # 获取文件扩展名
            file_ext = Path(file_name).suffix.lower()
            
            if file_ext not in self.supported_formats:
                raise DocumentProcessingException(f"不支持的文件格式: {file_ext}")
            
            # 在线程池中执行解析
            loop = asyncio.get_event_loop()
            result_data = await loop.run_in_executor(
                self.executor,
                self._parse_document_sync,
                file_path,
                file_name,
                file_ext,
                extract_images,
                extract_tables
            )
            
            # 创建解析结果
            processing_time = time.time() - start_time
            result = DocumentParseResult(
                content=result_data['content'],
                metadata=result_data['metadata'],
                images=result_data['images'],
                tables=result_data['tables'],
                structure=result_data['structure'],
                page_count=result_data['page_count'],
                word_count=len(result_data['content'].split()) if result_data['content'] else 0,
                file_hash=file_hash,
                processing_time=processing_time
            )
            
            # 缓存结果
            if use_cache:
                await self._cache_result(file_hash, result)
            
            logger.info(f"文档解析完成: {file_name}, 耗时: {processing_time:.2f}s")
            return result
            
        except Exception as e:
            logger.error(f"文档解析失败 {file_name}: {e}")
            raise DocumentProcessingException(f"文档解析失败: {e}")
    
    def _parse_document_sync(
        self,
        file_path: str,
        file_name: str,
        file_ext: str,
        extract_images: bool,
        extract_tables: bool
    ) -> Dict[str, Any]:
        """同步解析文档"""
        parser_func = self.supported_formats[file_ext]
        
        result = {
            'content': '',
            'metadata': {
                'file_name': file_name,
                'file_path': file_path,
                'file_type': file_ext,
                'file_size': os.path.getsize(file_path)
            },
            'images': [],
            'tables': [],
            'structure': {},
            'page_count': 0
        }
        
        # 调用对应的解析函数
        parsed_data = parser_func(file_path, extract_images, extract_tables)
        result.update(parsed_data)
        
        return result
    
    def _parse_pdf_with_marker(
        self, 
        file_path: str, 
        extract_images: bool, 
        extract_tables: bool
    ) -> Dict[str, Any]:
        """
        使用Marker解析PDF文档
        如果Marker不可用，降级到PyMuPDF
        """
        try:
            # 尝试使用Marker
            return self._parse_pdf_marker(file_path, extract_images, extract_tables)
        except ImportError:
            logger.warning("Marker未安装，使用PyMuPDF解析PDF")
            return self._parse_pdf_pymupdf(file_path, extract_images, extract_tables)
        except Exception as e:
            logger.warning(f"Marker解析失败，降级到PyMuPDF: {e}")
            return self._parse_pdf_pymupdf(file_path, extract_images, extract_tables)
    
    def _parse_pdf_marker(
        self, 
        file_path: str, 
        extract_images: bool, 
        extract_tables: bool
    ) -> Dict[str, Any]:
        """使用Marker解析PDF"""
        try:
            # 导入Marker（如果可用）
            from marker import convert_single_pdf
            from marker.models import load_all_models
            
            # 加载模型
            model_lst = load_all_models()
            
            # 转换PDF
            full_text, images, out_meta = convert_single_pdf(
                file_path,
                model_lst,
                max_pages=self.marker_config['max_pages'],
                langs=self.marker_config['languages'],
                batch_multiplier=self.marker_config['batch_size']
            )
            
            # 处理图片
            processed_images = []
            if extract_images and images:
                for i, img in enumerate(images):
                    try:
                        # 转换图片为base64
                        img_buffer = io.BytesIO()
                        img.save(img_buffer, format='PNG')
                        img_base64 = base64.b64encode(img_buffer.getvalue()).decode()
                        
                        processed_images.append({
                            'index': i,
                            'format': 'PNG',
                            'data': img_base64,
                            'size': img.size,
                            'source': 'marker'
                        })
                    except Exception as e:
                        logger.warning(f"处理图片失败 {i}: {e}")
            
            # 提取表格（从文本中识别）
            tables = []
            if extract_tables:
                tables = self._extract_tables_from_text(full_text)
            
            return {
                'content': full_text,
                'metadata': {
                    'parser': 'marker',
                    'marker_meta': out_meta,
                    'languages': self.marker_config['languages']
                },
                'images': processed_images,
                'tables': tables,
                'page_count': out_meta.get('pages', 0) if out_meta else 0,
                'structure': {
                    'type': 'pdf',
                    'parser': 'marker',
                    'has_images': len(processed_images) > 0,
                    'has_tables': len(tables) > 0
                }
            }
            
        except Exception as e:
            logger.error(f"Marker PDF解析失败: {e}")
            raise DocumentProcessingException(f"Marker PDF解析失败: {e}")
    
    def _parse_pdf_pymupdf(
        self, 
        file_path: str, 
        extract_images: bool, 
        extract_tables: bool
    ) -> Dict[str, Any]:
        """使用PyMuPDF解析PDF（降级方案）"""
        try:
            doc = fitz.open(file_path)
            content_parts = []
            images = []
            tables = []
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # 提取文本
                text = page.get_text()
                if text.strip():
                    content_parts.append(f"\n--- 第 {page_num + 1} 页 ---\n")
                    content_parts.append(text)
                
                # 提取图片
                if extract_images:
                    image_list = page.get_images()
                    for img_index, img in enumerate(image_list):
                        try:
                            xref = img[0]
                            pix = fitz.Pixmap(doc, xref)
                            
                            if pix.n - pix.alpha < 4:  # 确保是RGB或灰度图
                                img_data = pix.tobytes("png")
                                img_base64 = base64.b64encode(img_data).decode()
                                
                                images.append({
                                    'page': page_num + 1,
                                    'index': img_index,
                                    'format': 'PNG',
                                    'data': img_base64,
                                    'width': pix.width,
                                    'height': pix.height,
                                    'source': 'pymupdf'
                                })
                            
                            pix = None
                        except Exception as e:
                            logger.warning(f"提取图片失败 page {page_num + 1}, img {img_index}: {e}")
                
                # 提取表格
                if extract_tables:
                    try:
                        tables_on_page = page.find_tables()
                        for table_index, table in enumerate(tables_on_page):
                            table_data = table.extract()
                            tables.append({
                                'page': page_num + 1,
                                'index': table_index,
                                'data': table_data,
                                'bbox': table.bbox,
                                'source': 'pymupdf'
                            })
                    except Exception as e:
                        logger.warning(f"提取表格失败 page {page_num + 1}: {e}")
            
            doc.close()
            
            return {
                'content': '\n'.join(content_parts),
                'metadata': {
                    'parser': 'pymupdf',
                    'total_pages': len(doc)
                },
                'images': images,
                'tables': tables,
                'page_count': len(doc),
                'structure': {
                    'type': 'pdf',
                    'parser': 'pymupdf',
                    'has_images': len(images) > 0,
                    'has_tables': len(tables) > 0
                }
            }
            
        except Exception as e:
            logger.error(f"PyMuPDF PDF解析失败: {e}")
            raise DocumentProcessingException(f"PyMuPDF PDF解析失败: {e}")
    
    def _extract_tables_from_text(self, text: str) -> List[Dict[str, Any]]:
        """从文本中提取表格结构"""
        tables = []
        lines = text.split('\n')
        
        current_table = []
        in_table = False
        
        for line_num, line in enumerate(lines):
            # 简单的表格检测逻辑
            if '|' in line and line.count('|') >= 2:
                if not in_table:
                    in_table = True
                    current_table = []
                
                # 解析表格行
                cells = [cell.strip() for cell in line.split('|')]
                if cells and cells[0] == '':
                    cells = cells[1:]
                if cells and cells[-1] == '':
                    cells = cells[:-1]
                
                if cells:
                    current_table.append(cells)
            else:
                if in_table and current_table:
                    # 表格结束，保存表格
                    tables.append({
                        'index': len(tables),
                        'data': current_table,
                        'rows': len(current_table),
                        'cols': len(current_table[0]) if current_table else 0,
                        'line_start': line_num - len(current_table),
                        'line_end': line_num - 1,
                        'source': 'text_extraction'
                    })
                    current_table = []
                    in_table = False
        
        # 处理文档末尾的表格
        if in_table and current_table:
            tables.append({
                'index': len(tables),
                'data': current_table,
                'rows': len(current_table),
                'cols': len(current_table[0]) if current_table else 0,
                'line_start': len(lines) - len(current_table),
                'line_end': len(lines) - 1,
                'source': 'text_extraction'
            })
        
        return tables
    
    async def _calculate_file_hash(self, file_path: str) -> str:
        """计算文件哈希值"""
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()
    
    async def _get_cached_result(self, file_hash: str) -> Optional[DocumentParseResult]:
        """获取缓存的解析结果"""
        cache_file = self.cache_dir / f"{file_hash}.json"
        if cache_file.exists():
            try:
                with open(cache_file, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                return DocumentParseResult(**data)
            except Exception as e:
                logger.warning(f"读取缓存失败: {e}")
        return None
    
    async def _cache_result(self, file_hash: str, result: DocumentParseResult):
        """缓存解析结果"""
        cache_file = self.cache_dir / f"{file_hash}.json"
        try:
            with open(cache_file, 'w', encoding='utf-8') as f:
                json.dump(result.__dict__, f, ensure_ascii=False, indent=2)
        except Exception as e:
            logger.warning(f"缓存结果失败: {e}")


    def _parse_docx(
        self,
        file_path: str,
        extract_images: bool,
        extract_tables: bool
    ) -> Dict[str, Any]:
        """解析Word文档"""
        try:
            doc = DocxDocument(file_path)
            content_parts = []
            images = []
            tables = []

            # 提取段落内容
            for para in doc.paragraphs:
                if para.text.strip():
                    content_parts.append(para.text)

            # 提取表格
            if extract_tables:
                for table_idx, table in enumerate(doc.tables):
                    table_data = []
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text.strip())
                        table_data.append(row_data)

                    tables.append({
                        'index': table_idx,
                        'data': table_data,
                        'rows': len(table_data),
                        'cols': len(table_data[0]) if table_data else 0,
                        'source': 'python-docx'
                    })

            # 提取图片（基础实现）
            if extract_images:
                try:
                    from docx.document import Document as DocxDoc
                    from docx.oxml.table import CT_Tbl
                    from docx.oxml.text.paragraph import CT_P
                    from docx.table import _Cell, Table
                    from docx.text.paragraph import Paragraph

                    # 遍历文档中的所有关系，查找图片
                    for rel in doc.part.rels.values():
                        if "image" in rel.target_ref:
                            try:
                                img_data = rel.target_part.blob
                                img_base64 = base64.b64encode(img_data).decode()

                                images.append({
                                    'index': len(images),
                                    'format': rel.target_ref.split('.')[-1].upper(),
                                    'data': img_base64,
                                    'source': 'python-docx'
                                })
                            except Exception as e:
                                logger.warning(f"提取Word图片失败: {e}")
                except Exception as e:
                    logger.warning(f"Word图片提取功能不可用: {e}")

            return {
                'content': '\n\n'.join(content_parts),
                'metadata': {
                    'parser': 'python-docx',
                    'paragraphs': len(doc.paragraphs),
                    'tables_count': len(tables),
                    'images_count': len(images)
                },
                'images': images,
                'tables': tables,
                'page_count': 1,  # Word文档页数难以准确计算
                'structure': {
                    'type': 'docx',
                    'parser': 'python-docx',
                    'has_images': len(images) > 0,
                    'has_tables': len(tables) > 0
                }
            }

        except Exception as e:
            logger.error(f"Word文档解析失败: {e}")
            raise DocumentProcessingException(f"Word文档解析失败: {e}")

    def _parse_excel(
        self,
        file_path: str,
        extract_images: bool,
        extract_tables: bool
    ) -> Dict[str, Any]:
        """解析Excel文档"""
        try:
            # 使用pandas读取Excel
            excel_file = pd.ExcelFile(file_path)
            content_parts = []
            tables = []

            for sheet_idx, sheet_name in enumerate(excel_file.sheet_names):
                df = pd.read_excel(file_path, sheet_name=sheet_name)

                # 添加工作表标题
                content_parts.append(f"\n=== 工作表: {sheet_name} ===\n")

                # 转换为文本
                sheet_text = df.to_string(index=False)
                content_parts.append(sheet_text)

                # 保存表格数据
                if extract_tables and not df.empty:
                    table_data = [df.columns.tolist()] + df.values.tolist()
                    tables.append({
                        'sheet': sheet_name,
                        'index': sheet_idx,
                        'data': table_data,
                        'rows': len(table_data),
                        'cols': len(table_data[0]) if table_data else 0,
                        'source': 'pandas'
                    })

            return {
                'content': '\n'.join(content_parts),
                'metadata': {
                    'parser': 'pandas',
                    'sheets': excel_file.sheet_names,
                    'sheet_count': len(excel_file.sheet_names)
                },
                'images': [],  # Excel图片提取较复杂，暂不实现
                'tables': tables,
                'page_count': len(excel_file.sheet_names),
                'structure': {
                    'type': 'excel',
                    'parser': 'pandas',
                    'sheets': excel_file.sheet_names,
                    'has_tables': len(tables) > 0
                }
            }

        except Exception as e:
            logger.error(f"Excel文档解析失败: {e}")
            raise DocumentProcessingException(f"Excel文档解析失败: {e}")

    def _parse_pptx(
        self,
        file_path: str,
        extract_images: bool,
        extract_tables: bool
    ) -> Dict[str, Any]:
        """解析PowerPoint文档"""
        try:
            prs = Presentation(file_path)
            content_parts = []
            images = []
            tables = []

            for slide_idx, slide in enumerate(prs.slides):
                content_parts.append(f"\n=== 幻灯片 {slide_idx + 1} ===\n")

                # 提取文本
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_parts.append(shape.text)

                    # 提取表格
                    if extract_tables and shape.has_table:
                        table = shape.table
                        table_data = []
                        for row in table.rows:
                            row_data = []
                            for cell in row.cells:
                                row_data.append(cell.text.strip())
                            table_data.append(row_data)

                        tables.append({
                            'slide': slide_idx + 1,
                            'index': len(tables),
                            'data': table_data,
                            'rows': len(table_data),
                            'cols': len(table_data[0]) if table_data else 0,
                            'source': 'python-pptx'
                        })

            return {
                'content': '\n'.join(content_parts),
                'metadata': {
                    'parser': 'python-pptx',
                    'slides': len(prs.slides)
                },
                'images': images,
                'tables': tables,
                'page_count': len(prs.slides),
                'structure': {
                    'type': 'pptx',
                    'parser': 'python-pptx',
                    'slides': len(prs.slides),
                    'has_tables': len(tables) > 0
                }
            }

        except Exception as e:
            logger.error(f"PowerPoint文档解析失败: {e}")
            raise DocumentProcessingException(f"PowerPoint文档解析失败: {e}")

    def _parse_text(
        self,
        file_path: str,
        extract_images: bool,
        extract_tables: bool
    ) -> Dict[str, Any]:
        """解析纯文本文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

            # 从文本中提取表格
            tables = []
            if extract_tables:
                tables = self._extract_tables_from_text(content)

            return {
                'content': content,
                'metadata': {
                    'parser': 'text',
                    'encoding': 'utf-8',
                    'lines': len(content.split('\n'))
                },
                'images': [],
                'tables': tables,
                'page_count': 1,
                'structure': {
                    'type': 'text',
                    'parser': 'text',
                    'has_tables': len(tables) > 0
                }
            }

        except UnicodeDecodeError:
            # 尝试其他编码
            for encoding in ['gbk', 'gb2312', 'latin1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()

                    tables = []
                    if extract_tables:
                        tables = self._extract_tables_from_text(content)

                    return {
                        'content': content,
                        'metadata': {
                            'parser': 'text',
                            'encoding': encoding,
                            'lines': len(content.split('\n'))
                        },
                        'images': [],
                        'tables': tables,
                        'page_count': 1,
                        'structure': {
                            'type': 'text',
                            'parser': 'text',
                            'has_tables': len(tables) > 0
                        }
                    }
                except UnicodeDecodeError:
                    continue

            raise DocumentProcessingException("无法识别文本文件编码")

        except Exception as e:
            logger.error(f"文本文件解析失败: {e}")
            raise DocumentProcessingException(f"文本文件解析失败: {e}")


    def _parse_markdown(
        self,
        file_path: str,
        extract_images: bool,
        extract_tables: bool
    ) -> Dict[str, Any]:
        """解析Markdown文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

            # 提取表格（Markdown表格格式）
            tables = []
            if extract_tables:
                import re
                table_pattern = r'\|.*\|'
                lines = content.split('\n')
                current_table = []
                in_table = False

                for line in lines:
                    if re.match(table_pattern, line.strip()):
                        if not in_table:
                            in_table = True
                            current_table = []

                        # 跳过分隔行
                        if not re.match(r'\|[\s\-\|:]+\|', line.strip()):
                            cells = [cell.strip() for cell in line.split('|')[1:-1]]
                            current_table.append(cells)
                    else:
                        if in_table and current_table:
                            tables.append({
                                'index': len(tables),
                                'data': current_table,
                                'rows': len(current_table),
                                'cols': len(current_table[0]) if current_table else 0,
                                'source': 'markdown'
                            })
                            current_table = []
                            in_table = False

            return {
                'content': content,
                'metadata': {
                    'parser': 'markdown',
                    'encoding': 'utf-8'
                },
                'images': [],
                'tables': tables,
                'page_count': 1,
                'structure': {
                    'type': 'markdown',
                    'parser': 'markdown',
                    'has_tables': len(tables) > 0
                }
            }

        except Exception as e:
            logger.error(f"Markdown文件解析失败: {e}")
            raise DocumentProcessingException(f"Markdown文件解析失败: {e}")

    def _parse_html(
        self,
        file_path: str,
        extract_images: bool,
        extract_tables: bool
    ) -> Dict[str, Any]:
        """解析HTML文件"""
        try:
            from bs4 import BeautifulSoup

            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()

            soup = BeautifulSoup(html_content, 'html.parser')

            # 提取文本内容
            text_content = soup.get_text(separator='\n', strip=True)

            # 提取表格
            tables = []
            if extract_tables:
                html_tables = soup.find_all('table')
                for table_idx, table in enumerate(html_tables):
                    table_data = []
                    rows = table.find_all('tr')
                    for row in rows:
                        cells = row.find_all(['td', 'th'])
                        row_data = [cell.get_text(strip=True) for cell in cells]
                        if row_data:
                            table_data.append(row_data)

                    if table_data:
                        tables.append({
                            'index': table_idx,
                            'data': table_data,
                            'rows': len(table_data),
                            'cols': len(table_data[0]) if table_data else 0,
                            'source': 'html'
                        })

            # 提取元数据
            metadata = {'parser': 'beautifulsoup'}
            title_tag = soup.find('title')
            if title_tag:
                metadata['title'] = title_tag.get_text()

            return {
                'content': text_content,
                'metadata': metadata,
                'images': [],
                'tables': tables,
                'page_count': 1,
                'structure': {
                    'type': 'html',
                    'parser': 'beautifulsoup',
                    'has_tables': len(tables) > 0
                }
            }

        except Exception as e:
            logger.error(f"HTML文件解析失败: {e}")
            raise DocumentProcessingException(f"HTML文件解析失败: {e}")

    def _parse_csv(
        self,
        file_path: str,
        extract_images: bool,
        extract_tables: bool
    ) -> Dict[str, Any]:
        """解析CSV文件"""
        try:
            df = pd.read_csv(file_path)

            # 转换为文本
            content = df.to_string(index=False)

            # CSV本身就是表格
            tables = []
            if extract_tables and not df.empty:
                table_data = [df.columns.tolist()] + df.values.tolist()
                tables.append({
                    'index': 0,
                    'data': table_data,
                    'rows': len(table_data),
                    'cols': len(table_data[0]) if table_data else 0,
                    'source': 'csv'
                })

            return {
                'content': content,
                'metadata': {
                    'parser': 'pandas',
                    'rows': len(df),
                    'columns': len(df.columns)
                },
                'images': [],
                'tables': tables,
                'page_count': 1,
                'structure': {
                    'type': 'csv',
                    'parser': 'pandas',
                    'has_tables': len(tables) > 0
                }
            }

        except Exception as e:
            logger.error(f"CSV文件解析失败: {e}")
            raise DocumentProcessingException(f"CSV文件解析失败: {e}")

    def _parse_json(
        self,
        file_path: str,
        extract_images: bool,
        extract_tables: bool
    ) -> Dict[str, Any]:
        """解析JSON文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)

            # 格式化JSON为可读文本
            content = json.dumps(data, ensure_ascii=False, indent=2)

            return {
                'content': content,
                'metadata': {
                    'parser': 'json',
                    'data_type': type(data).__name__
                },
                'images': [],
                'tables': [],
                'page_count': 1,
                'structure': {
                    'type': 'json',
                    'parser': 'json',
                    'data_structure': type(data).__name__
                }
            }

        except Exception as e:
            logger.error(f"JSON文件解析失败: {e}")
            raise DocumentProcessingException(f"JSON文件解析失败: {e}")

    def _parse_xml(
        self,
        file_path: str,
        extract_images: bool,
        extract_tables: bool
    ) -> Dict[str, Any]:
        """解析XML文件"""
        try:
            import xml.etree.ElementTree as ET
            from bs4 import BeautifulSoup

            with open(file_path, 'r', encoding='utf-8') as f:
                xml_content = f.read()

            # 使用BeautifulSoup解析XML并提取文本
            soup = BeautifulSoup(xml_content, 'xml')
            text_content = soup.get_text(separator='\n', strip=True)

            return {
                'content': text_content,
                'metadata': {
                    'parser': 'xml',
                    'encoding': 'utf-8'
                },
                'images': [],
                'tables': [],
                'page_count': 1,
                'structure': {
                    'type': 'xml',
                    'parser': 'xml'
                }
            }

        except Exception as e:
            logger.error(f"XML文件解析失败: {e}")
            raise DocumentProcessingException(f"XML文件解析失败: {e}")

    def _parse_rtf(
        self,
        file_path: str,
        extract_images: bool,
        extract_tables: bool
    ) -> Dict[str, Any]:
        """解析RTF文件"""
        try:
            # RTF解析较复杂，这里提供基础实现
            with open(file_path, 'r', encoding='utf-8') as f:
                rtf_content = f.read()

            # 简单的RTF文本提取（移除控制字符）
            import re
            # 移除RTF控制字符
            text = re.sub(r'\\[a-z]+\d*\s?', '', rtf_content)
            text = re.sub(r'[{}]', '', text)
            text = text.strip()

            return {
                'content': text,
                'metadata': {
                    'parser': 'rtf_simple',
                    'encoding': 'utf-8'
                },
                'images': [],
                'tables': [],
                'page_count': 1,
                'structure': {
                    'type': 'rtf',
                    'parser': 'rtf_simple'
                }
            }

        except Exception as e:
            logger.error(f"RTF文件解析失败: {e}")
            raise DocumentProcessingException(f"RTF文件解析失败: {e}")


# 全局Marker文档解析器实例
marker_parser = MarkerDocumentParser()
